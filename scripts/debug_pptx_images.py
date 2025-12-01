import argparse
import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

def analyze_pptx_images(pptx_path):
    print(f"Analyzing PPTX: {pptx_path}")
    
    if not pptx_path.exists():
        print("Error: File not found!")
        return

    prs = Presentation(pptx_path)
    
    total_images = 0
    valid_images = 0
    corrupt_images = 0
    unsupported_formats = {}

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_images += 1
                try:
                    image = shape.image
                    content_type = image.content_type
                    ext = image.ext
                    blob = image.blob
                    
                    print(f"\n[Slide {i+1}] Shape ID {shape.shape_id}")
                    print(f"  - Content Type: {content_type}")
                    print(f"  - Extension: {ext}")
                    print(f"  - Size: {len(blob)} bytes")
                    
                    # Try to open with Pillow
                    try:
                        with Image.open(io.BytesIO(blob)) as img:
                            print(f"  - Pillow Check: OK (Format: {img.format}, Size: {img.size}, Mode: {img.mode})")
                            valid_images += 1
                    except Exception as e:
                        print(f"  - Pillow Check: FAILED ({e})")
                        corrupt_images += 1
                        
                        # Count failure types
                        key = f"{ext} ({content_type})"
                        unsupported_formats[key] = unsupported_formats.get(key, 0) + 1

                except Exception as e:
                    print(f"  - Error accessing image property: {e}")

    print("\n" + "="*30)
    print("SUMMARY")
    print("="*30)
    print(f"Total Pictures Found: {total_images}")
    print(f"Valid (Openable by Pillow): {valid_images}")
    print(f"Corrupt/Unsupported: {corrupt_images}")
    
    if unsupported_formats:
        print("\nBreakdown of Unsupported Formats:")
        for fmt, count in unsupported_formats.items():
            print(f"  - {fmt}: {count}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("file_path", type=str, help="Path to the PPTX file")
    args = parser.parse_args()
    
    analyze_pptx_images(Path(args.file_path))
