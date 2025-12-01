import os
import json
import argparse
from pathlib import Path
from typing import List, Dict, Tuple
from tqdm import tqdm
import docx
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def save_image(image_bytes: bytes, filename: str, output_dir: Path) -> str:
    """Saves image bytes to disk and returns the path."""
    if not output_dir.exists():
        output_dir.mkdir(parents=True, exist_ok=True)
    
    file_path = output_dir / filename
    with open(file_path, "wb") as f:
        f.write(image_bytes)
    return str(file_path)

def extract_from_docx(file_path: Path, img_output_dir: Path) -> Tuple[str, List[str]]:
    """Extract text and images from DOCX."""
    text_content = []
    extracted_images = []
    
    try:
        doc = docx.Document(file_path)
        
        # Text Extraction
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text.strip())
        
        for table in doc.tables:
            text_content.append("\n--- Table ---")
            for row in table.rows:
                row_text = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                text_content.append(" | ".join(row_text))
            text_content.append("--- End Table ---\n")

        # Image Extraction (Basic support for 'rId' relationships)
        # Note: python-docx doesn't have a high-level API for image extraction order relative to text easily.
        # We will extract all images and append references at the end or try to find blips.
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_bytes = rel.target_part.blob
                    # Attempt to guess extension
                    content_type = rel.target_part.content_type
                    ext = ".jpg"
                    if "png" in content_type: ext = ".png"
                    elif "jpeg" in content_type: ext = ".jpg"
                    
                    img_filename = f"{file_path.stem}_img_{len(extracted_images)}{ext}"
                    saved_path = save_image(image_bytes, img_filename, img_output_dir)
                    extracted_images.append(saved_path)
                    
                    # Add a placeholder marker
                    text_content.append(f"\n[IMAGE_REF: {saved_path}]\n")
                except Exception as e:
                    pass # Skip unreadable images

        return "\n".join(text_content), extracted_images

    except Exception as e:
        print(f"Error reading DOCX {file_path.name}: {e}")
        return "", []

def extract_from_pdf(file_path: Path, img_output_dir: Path) -> Tuple[str, List[str]]:
    """Extract text and images from PDF."""
    text_content = []
    extracted_images = []
    
    try:
        reader = PdfReader(file_path)
        for page_num, page in enumerate(reader.pages):
            # Text
            extracted_text = page.extract_text()
            if extracted_text:
                text_content.append(extracted_text)
            
            # Images
            if "/XObject" in page["/Resources"]:
                xObject = page["/Resources"]["/XObject"].get_object()
                for obj in xObject:
                    if xObject[obj]["/Subtype"] == "/Image":
                        try:
                            image_obj = xObject[obj]
                            data = image_obj.get_data()
                            
                            # Filter small icons/lines
                            if len(data) < 1024: continue 
                            
                            ext = ".jpg" # Default
                            if "/Filter" in image_obj:
                                if "/FlateDecode" in image_obj["/Filter"]:
                                    ext = ".png"
                                elif "/DCTDecode" in image_obj["/Filter"]:
                                    ext = ".jpg"

                            img_filename = f"{file_path.stem}_page{page_num+1}_{obj[1:]}{ext}"
                            saved_path = save_image(data, img_filename, img_output_dir)
                            extracted_images.append(saved_path)
                            
                            text_content.append(f"\n[IMAGE_REF: {saved_path}]\n")
                        except Exception:
                            continue

        return "\n".join(text_content), extracted_images
    except Exception as e:
        print(f"Error reading PDF {file_path.name}: {e}")
        return "", []

def extract_from_pptx(file_path: Path, img_output_dir: Path) -> Tuple[str, List[str]]:
    """Extract text and images from PPTX."""
    text_content = []
    extracted_images = []
    
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            text_content.append(f"\n--- Slide {i+1} ---")
            
            for shape in slide.shapes:
                # Text
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
                
                # Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = "." + image.ext
                        
                        img_filename = f"{file_path.stem}_slide{i+1}_shape{shape.shape_id}{ext}"
                        saved_path = save_image(image_bytes, img_filename, img_output_dir)
                        extracted_images.append(saved_path)
                        
                        slide_text.append(f"\n[IMAGE_REF: {saved_path}]\n")
                    except Exception:
                        pass

            if slide_text:
                text_content.extend(slide_text)
                
        return "\n".join(text_content), extracted_images
    except Exception as e:
        print(f"Error reading PPTX {file_path.name}: {e}")
        return "", []

def extract_from_text_file(file_path: Path, img_output_dir: Path) -> Tuple[str, List[str]]:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(), []
    except Exception as e:
        return "", []

def clean_text(text: str) -> str:
    import re
    # Collapse multiple newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

def main():
    parser = argparse.ArgumentParser(description="Extract raw text AND images from Product Manuals")
    parser.add_argument("--input_dir", type=str, required=True, help="Directory containing raw files")
    parser.add_argument("--output_file", type=str, required=True, help="Output JSONL file")
    parser.add_argument("--image_output_dir", type=str, default="./extracted_assets", help="Dir to save extracted images")
    args = parser.parse_args()

    input_path = Path(args.input_dir).expanduser().resolve()
    img_out_path = Path(args.image_output_dir).resolve()
    
    results = []
    all_extracted_images = []

    print(f"Scanning {input_path}...")
    print(f"Images will be saved to {img_out_path}...")
    
    supported_extensions = {
        '.docx': extract_from_docx,
        '.pdf': extract_from_pdf,
        '.pptx': extract_from_pptx,
        '.md': extract_from_text_file,
        '.txt': extract_from_text_file
    }
    
    files_to_process = []
    for root, _, files in os.walk(input_path):
        for file in files:
            ext = Path(file).suffix.lower()
            if ext in supported_extensions:
                files_to_process.append(Path(root) / file)

    print(f"Found {len(files_to_process)} documents.")

    total_chars = 0
    for file_path in tqdm(files_to_process, desc="Extracting"):
        ext = file_path.suffix.lower()
        extractor = supported_extensions[ext]
        
        text, imgs = extractor(file_path, img_out_path)
        
        if imgs:
            all_extracted_images.extend(imgs)
        
        if text:
            text = clean_text(text)
            if len(text) > 20: # Lower threshold
                total_chars += len(text)
                results.append({
                    "file_path": str(file_path),
                    "filename": file_path.name,
                    "content": text,
                    "product_area": file_path.parent.name,
                    "extracted_images": imgs # Keep track of associated images
                })

    print(f"\n=== Extraction Summary ===")
    print(f"Processed Files: {len(results)}")
    print(f"Total Characters: {total_chars}")
    print(f"Extracted Images: {len(all_extracted_images)}")
    
    with open(args.output_file, 'w', encoding='utf-8') as f:
        for entry in results:
            f.write(json.dumps(entry, ensure_ascii=False) + "\n")
    
    print(f"Saved extracted text data to {args.output_file}")

if __name__ == "__main__":
    main()
